from dotenv import load_dotenv
import os
from flask import Flask, render_template, request, redirect, url_for, flash
from pptx import Presentation
import openai
import json
import jsonschema
from jsonschema import validate
import PyPDF2

load_dotenv()

client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


app = Flask(__name__)
app.secret_key = 'sushanttiwari098765432123456789'

schema_file_path = os.path.join(os.path.dirname(__file__), 'nclex_schema.json')

with open(schema_file_path, 'r') as f:
    NCLEX_SCHEMA = json.load(f)

def validate_questions(json_data):
    try:
        validate(instance=json_data, schema=NCLEX_SCHEMA)
        return True, None
    except jsonschema.exceptions.ValidationError as err:
        return False, err

UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_runs.append(text)
    return "\n".join(text_runs)

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text.strip()


def build_quiz_html(questions_json):
    quiz_html = '<form id="quiz-form">\n'
    
    for idx, question in enumerate(questions_json.get("questions", []), start=1):
        html_question = f'''
<div class="quiz-box" id="quiz-{idx}">
<div class="mt-3 fw-bold" id="score-output"></div>
  <h2>{idx}. {question["stem"]}</h2>
  <div class="answers">
'''
        input_type = "checkbox" if isinstance(question["correctAnswer"], list) else "radio"
        for option in question["options"]:
            input_id = f"q{idx}{option['label'].lower()}"
            input_name = f"question-{idx}" + ("[]" if input_type == "checkbox" else "")
            html_question += f'''
    <input type="{input_type}" id="{input_id}" name="{input_name}" value="{option["label"]}">
    <label for="{input_id}">{option["label"]}. {option["text"]}</label>
'''
        
        html_question += f'''
  </div>
  <span id="correct-answer-{idx}" style="display:none;">{question["correctAnswer"]}</span>
  <div class="explanation" id="explanation-{idx}" data-rationale="{question["explanation"]}" style="display:none;"></div>
</div>
'''
        quiz_html += html_question

    quiz_html += '''
  <div class="text-center mt-4">
    <button type="button" class="btn btn-success" onclick="gradeQuiz()">Submit Quiz</button>
    <div class="mt-3 fw-bold" id="score-output"></div>
  </div>
</form>
'''
    return quiz_html



def generate_nclex_questions(content, examples, num_questions, include_sata=False):
    sata_instruction = """
You MUST include at least 2 Select-All-That-Apply (SATA) questions. SATA questions must have exactly 5 options and use a list of letters for the correctAnswer field (e.g., ["A", "C", "E"]). Regular questions should have 4 options with a single correct answer (e.g., "B").
""" if include_sata else """
All questions must be single-answer multiple-choice (no SATA). Use exactly 4 options with a single correct answer (e.g., "B").
"""

    prompt = f"""
You are an expert nurse educator. Based ONLY on the following nursing content, generate exactly {num_questions} NCLEX-style questions.

{sata_instruction}

Each question MUST follow this strict JSON format with NO duplicated options and NO repeated answer lists:
{{
  "questionID": "string (unique ID like 101, 102...)",
  "stem": "A single clinical scenario question. Do not repeat options.",
  "options": [
    {{"label": "A", "text": "Option text"}},
    {{"label": "B", "text": "Option text"}},
    {{"label": "C", "text": "Option text"}},
    {{"label": "D", "text": "Option text"}}{"," if include_sata else ""}
    {{"label": "E", "text": "Option text"}}  // Only include this for SATA questions
  ],
  "correctAnswer": "A" | "B" | "C" | "D"{' OR ["A", "C", "E"]  // Only for SATA' if include_sata else ""},
  "explanation": "Brief, accurate rationale explaining why this is the correct answer.",
  "tags": ["topic1", "topic2"]
}}

Important instructions:
- Do NOT include commentary, markdown, or any text before or after the JSON.
- Only output a VALID JSON object with a top-level "questions" list.
- All questions must resemble the tone, logic, and structure of these examples:
{examples}

Nursing Content:
{content}
"""

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant specialized in generating NCLEX-style exam questions."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.5,
        max_tokens=3000
    )

    return response.choices[0].message.content


def load_examples(file_path):
    with open(file_path, 'r') as f:
        data = json.load(f)
    return json.dumps(data["examples"], indent=2)


examples_file_path = os.path.join(os.path.dirname(__file__), 'examples.json')
examples_content = load_examples(examples_file_path)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        if "pptx_file" not in request.files:
            flash("No file part")
            return redirect(request.url)
        file = request.files["pptx_file"]
        if file.filename == "":
            flash("No selected file")
            return redirect(request.url)
        if file:
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)
    
            if file.filename.lower().endswith('.pptx'):
                extracted_text = extract_text_from_ppt(file_path)
            elif file.filename.lower().endswith('.pdf'):
                extracted_text = extract_text_from_pdf(file_path)
            else:
                flash("Unsupported file type. Please upload a .pptx or .pdf file.")
                return redirect(request.url)
    
            if not extracted_text:
                flash("No text found in the file.")
                return redirect(request.url)
            
            try:
                num_questions = int(request.form.get("num_questions", 1))
            except ValueError:
                num_questions = 1
                
            include_sata = request.form.get("include_sata") == "true"

            try:
                raw_output = generate_nclex_questions(
                    extracted_text,
                    examples_content,
                    num_questions,
                    include_sata=include_sata
                )

                questions_json = json.loads(raw_output) 
                quiz_html = build_quiz_html(questions_json)
            except Exception as e:
                flash(f"Error processing generated questions: {str(e)}")
                return redirect(request.url)

    
            return render_template("index.html", quiz=quiz_html)
    
    return render_template("index.html")


if __name__ == "__main__":
    app.run(debug=True)
